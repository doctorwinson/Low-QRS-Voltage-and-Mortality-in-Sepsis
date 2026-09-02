#!/usr/bin/env python3
"""Generate SCI-grade fixed-coordinate cohort selection flowcharts.

Input: JSON by default. YAML is supported if PyYAML is installed.
Output: PDF, SVG, 600-dpi PNG, editable PPTX when python-pptx is available,
and a QC report.
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import textwrap
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, Rectangle

try:
    import yaml
except ImportError:  # pragma: no cover
    yaml = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt
except ImportError:  # pragma: no cover
    Presentation = None


DEFAULT_CONFIG: dict[str, Any] = {
    "figure_title": "Figure 1. Cohort selection flow diagram for the 72-hour landmark analysis",
    "cohorts": [
        {
            "name": "MIMIC-IV cohort",
            "main_flow": [
                {"label": "Source sepsis cohort", "n": "25,311"},
                {"label": "SIC-core variables reconstructable", "n": "18,154"},
                {"label": "Free of SIC through the 72-h landmark", "n": "6,660"},
                {"label": "Alive at the 72-h landmark", "n": "6,550"},
                {"label": "72-h landmark risk set", "n": "4,739"},
            ],
            "exclusions": [
                {"after_step": 1, "label": "Excluded: insufficient data for SIC-core reconstruction", "n": "7,157"},
                {"after_step": 2, "label": "Excluded: SIC at or before 72 h", "n": "11,494"},
                {"after_step": 3, "label": "Excluded: death at or before 72 h", "n": "110"},
                {
                    "after_step": 4,
                    "label": "Excluded: discharged before the landmark or unobservable after 72 h",
                    "n": "1,811",
                },
            ],
            "final_groups": [
                {"label": "LP-SIC events after the 72-h landmark", "n": "274"},
                {"label": "No LP-SIC event", "n": "4,465"},
            ],
        },
        {
            "name": "eICU-CRD cohort",
            "main_flow": [
                {"label": "Source sepsis cohort", "n": "20,785"},
                {"label": "SIC-core variables reconstructable", "n": "4,892"},
                {"label": "Free of SIC through the 72-h landmark", "n": "1,127"},
                {"label": "Alive at the 72-h landmark", "n": "1,119"},
                {"label": "72-h landmark risk set", "n": "801"},
            ],
            "exclusions": [
                {"after_step": 1, "label": "Excluded: insufficient data for SIC-core reconstruction", "n": "15,893"},
                {"after_step": 2, "label": "Excluded: SIC at or before 72 h", "n": "3,765"},
                {"after_step": 3, "label": "Excluded: death at or before 72 h", "n": "8"},
                {
                    "after_step": 4,
                    "label": "Excluded: discharged before the landmark or unobservable after 72 h",
                    "n": "318",
                },
            ],
            "final_groups": [
                {"label": "LP-SIC events after the 72-h landmark", "n": "68"},
                {"label": "No LP-SIC event", "n": "733"},
            ],
        },
    ],
}


@dataclass(frozen=True)
class Box:
    key: str
    cx: float
    cy: float
    w: float
    h: float
    text: str
    kind: str

    @property
    def left(self) -> float:
        return self.cx - self.w / 2

    @property
    def right(self) -> float:
        return self.cx + self.w / 2

    @property
    def bottom(self) -> float:
        return self.cy - self.h / 2

    @property
    def top(self) -> float:
        return self.cy + self.h / 2

    @property
    def top_center(self) -> tuple[float, float]:
        return (self.cx, self.top)

    @property
    def bottom_center(self) -> tuple[float, float]:
        return (self.cx, self.bottom)

    @property
    def left_center(self) -> tuple[float, float]:
        return (self.left, self.cy)

    @property
    def right_center(self) -> tuple[float, float]:
        return (self.right, self.cy)


@dataclass(frozen=True)
class OutcomeSplit:
    trunk_start: tuple[float, float]
    junction: tuple[float, float]
    branch_junctions: tuple[tuple[float, float], ...]
    child_ends: tuple[tuple[float, float], ...]


LAYOUT_TOL = 1e-6


def parse_n(value: Any) -> int:
    if isinstance(value, int):
        return value
    text = str(value)
    digits = re.sub(r"[^0-9-]", "", text)
    if not digits:
        raise ValueError(f"Cannot parse n from {value!r}")
    return int(digits)


def format_n(value: Any) -> str:
    return f"{parse_n(value):,}"


def wrap_label(label: str, n: Any | None = None, width: int = 26) -> str:
    clean = str(label).strip()
    clean = re.sub(r"^Excluded:\s*", "Excluded: ", clean, flags=re.IGNORECASE)
    wrapped = textwrap.fill(clean, width=width, break_long_words=False)
    if n is None or str(n).strip() == "":
        return wrapped
    return f"{wrapped}\nn = {format_n(n)}"


def load_config(path: Path | None) -> dict[str, Any]:
    if path is None:
        return DEFAULT_CONFIG
    text = path.read_text(encoding="utf-8-sig")
    if path.suffix.lower() in {".yaml", ".yml"}:
        if yaml is None:
            raise RuntimeError("YAML input requires PyYAML. Use JSON or install PyYAML.")
        return yaml.safe_load(text)
    return json.loads(text)


def layout_float(config: dict[str, Any], key: str, default: float) -> float:
    layout = config.get("layout", {})
    if not isinstance(layout, dict):
        return default
    return float(layout.get(key, default))


def outcome_centers(main_x: float, outcome_w: float, outcome_gap: float, group_count: int) -> list[float]:
    if not 2 <= group_count <= 5:
        raise ValueError("This template supports 2 to 5 final groups per cohort.")
    if outcome_gap < 0:
        raise ValueError("layout.outcome_gap must be non-negative.")
    total_width = outcome_w * group_count + outcome_gap * (group_count - 1)
    first_center = main_x - total_width / 2 + outcome_w / 2
    return [first_center + i * (outcome_w + outcome_gap) for i in range(group_count)]


def outcome_split_geometry(risk_box: Box, group_boxes: list[Box]) -> OutcomeSplit:
    if not 2 <= len(group_boxes) <= 5:
        raise ValueError("Outcome split geometry requires 2 to 5 final group boxes.")
    child_top = max(box.top for box in group_boxes)
    if risk_box.bottom <= child_top:
        raise ValueError(
            f"Insufficient vertical space for outcome split under {risk_box.key}: "
            f"risk bottom={risk_box.bottom:.3f}, child top={child_top:.3f}."
        )
    y_joint = (risk_box.bottom + child_top) / 2
    junction = (risk_box.cx, y_joint)
    return OutcomeSplit(
        trunk_start=risk_box.bottom_center,
        junction=junction,
        branch_junctions=tuple((box.cx, y_joint) for box in group_boxes),
        child_ends=tuple(box.top_center for box in group_boxes),
    )


def is_layout_pass(lines: list[str]) -> bool:
    return not any(line.startswith("- FAIL:") for line in lines)


def layout_params(config: dict[str, Any]) -> dict[str, float]:
    return {
        "main_w": layout_float(config, "main_w", 2.35),
        "main_h": layout_float(config, "main_h", 0.65),
        "excl_w": layout_float(config, "exclusion_w", 2.75),
        "excl_h": layout_float(config, "exclusion_h", 0.62),
        "exclusion_gap": layout_float(config, "exclusion_gap", 0.65),
        "outcome_w": layout_float(config, "outcome_w", 1.65),
        "outcome_h": layout_float(config, "outcome_h", 0.65),
        "outcome_gap": layout_float(config, "outcome_gap", 0.25),
        "outcome_visible_gap": layout_float(config, "outcome_visible_gap", 0.65),
        "panel_gap": layout_float(config, "panel_gap", 0.85),
        "panel_margin": layout_float(config, "panel_margin", 0.50),
        "main_y_start": layout_float(config, "main_y_start", 8.65),
        "main_step_gap": layout_float(config, "main_step_gap", 1.30),
    }


def outcome_total_width(group_count: int, outcome_w: float, outcome_gap: float) -> float:
    if group_count <= 0:
        return 0.0
    return outcome_w * group_count + outcome_gap * max(group_count - 1, 0)


def compute_x_layouts(config: dict[str, Any]) -> tuple[float, list[dict[str, float]]]:
    cohorts = config["cohorts"]
    if len(cohorts) > 2:
        raise ValueError("This fixed template supports one or two cohorts. Split larger figures into panels.")

    params = layout_params(config)
    layouts: list[dict[str, float]] = []
    x_cursor = 0.0

    for cohort in cohorts:
        group_count = len(cohort.get("final_groups", []))
        if group_count and not 2 <= group_count <= 5:
            raise ValueError("Each cohort must have either no final groups or 2 to 5 final groups.")
        effective_group_count = max(group_count, 2)
        total_outcome = outcome_total_width(effective_group_count, params["outcome_w"], params["outcome_gap"])
        left_extent = max(params["main_w"] / 2, total_outcome / 2)
        right_extent = max(
            params["main_w"] / 2 + params["exclusion_gap"] + params["excl_w"],
            total_outcome / 2,
        )
        panel_w = left_extent + right_extent + 2 * params["panel_margin"]
        main_x = x_cursor + params["panel_margin"] + left_extent
        excl_x = main_x + params["main_w"] / 2 + params["exclusion_gap"] + params["excl_w"] / 2
        layouts.append({"main_x": main_x, "excl_x": excl_x, "panel_w": panel_w})
        x_cursor += panel_w + params["panel_gap"]

    computed_fig_w = x_cursor - params["panel_gap"] if layouts else 8.0
    requested_fig_w = layout_float(config, "fig_w", computed_fig_w)
    return max(computed_fig_w, requested_fig_w, 8.0), layouts


def validate_counts(config: dict[str, Any]) -> tuple[bool, list[str]]:
    if str(config.get("diagram_mode", "cohort")).lower() == "roadmap":
        lines = [
            "## Workflow Consistency",
            "- PASS: roadmap mode omits sample-size arithmetic for methodological steps.",
        ]
        ok = True
        for cohort in config["cohorts"]:
            main = cohort.get("main_flow", [])
            if not main:
                ok = False
                lines.append(f"- FAIL: {cohort.get('name', 'Unnamed workflow')} has no main-flow steps.")
                continue
            for item in cohort.get("exclusions", []):
                step = int(item["after_step"])
                if not 1 <= step < len(main):
                    ok = False
                    lines.append(f"- FAIL: side note after_step={step} is outside the workflow transitions.")
            group_count = len(cohort.get("final_groups", []))
            if group_count and not 2 <= group_count <= 5:
                ok = False
                lines.append(f"- FAIL: roadmap output branches must contain 2 to 5 groups; found {group_count}.")
        if ok:
            lines.append("- PASS: all side notes and output branches map to valid workflow anchors.")
        return ok, lines

    ok = True
    lines = ["## Sample-Size Consistency"]
    for cohort in config["cohorts"]:
        name = cohort["name"]
        main = cohort["main_flow"]
        exclusions = {int(item["after_step"]): item for item in cohort.get("exclusions", [])}
        lines.append(f"### {name}")
        expected_transitions = max(len(main) - 1, 0)
        if len(exclusions) != expected_transitions:
            ok = False
            lines.append(
                f"- FAIL: {len(exclusions)} exclusions provided for {expected_transitions} main-flow transitions."
            )
        for step in range(1, len(main)):
            prev_n = parse_n(main[step - 1]["n"])
            next_n = parse_n(main[step]["n"])
            exclusion = exclusions.get(step)
            if exclusion is None:
                ok = False
                lines.append(f"- FAIL: missing exclusion after step {step}.")
                continue
            excluded_n = parse_n(exclusion["n"])
            expected = prev_n - excluded_n
            if expected == next_n:
                lines.append(f"- PASS: step {step}: {prev_n:,} - {excluded_n:,} = {next_n:,}.")
            else:
                ok = False
                lines.append(
                    f"- FAIL: step {step}: previous n = {prev_n:,}, excluded n = {excluded_n:,}, "
                    f"expected n = {expected:,}, provided n = {next_n:,}."
                )
        groups = cohort.get("final_groups", [])
        if groups:
            final_n = parse_n(main[-1]["n"])
            group_sum = sum(parse_n(group["n"]) for group in groups)
            if final_n == group_sum:
                lines.append(f"- PASS: final groups sum to final risk set n = {final_n:,}.")
            else:
                ok = False
                lines.append(
                    f"- FAIL: final group sum = {group_sum:,}; final risk set n = {final_n:,}."
                )
    return ok, lines


def make_layout(config: dict[str, Any]) -> tuple[dict[str, Box], list[str]]:
    cohorts = config["cohorts"]
    if len(cohorts) > 2:
        raise ValueError("This fixed template supports one or two cohorts. Split larger figures into panels.")

    params = layout_params(config)
    main_w = params["main_w"]
    main_h = params["main_h"]
    excl_w = params["excl_w"]
    excl_h = params["excl_h"]
    outcome_w = params["outcome_w"]
    outcome_h = params["outcome_h"]
    outcome_gap = params["outcome_gap"]
    outcome_visible_gap = params["outcome_visible_gap"]
    _, x_layouts = compute_x_layouts(config)

    boxes: dict[str, Box] = {}
    lines = ["## Layout Consistency"]
    main_lengths = {len(cohort["main_flow"]) for cohort in cohorts}
    if len(main_lengths) == 1:
        lines.append("- PASS: all cohorts share the same number of main-flow steps.")
    else:
        lines.append("- FAIL: cohorts have different numbers of main-flow steps; y alignment needs manual review.")

    max_steps = max(main_lengths)
    if max_steps > 6:
        raise ValueError("Template currently supports up to 6 main-flow steps.")
    y_main = [
        params["main_y_start"] - idx * params["main_step_gap"]
        for idx in range(max_steps)
    ]

    y_outcome = y_main[max_steps - 1] - main_h / 2 - outcome_visible_gap - outcome_h / 2

    for cohort_idx, cohort in enumerate(cohorts):
        layout = x_layouts[cohort_idx]
        prefix = f"cohort_{cohort_idx}"
        boxes[f"{prefix}_title"] = Box(
            f"{prefix}_title",
            layout["main_x"],
            9.35,
            main_w,
            0.35,
            str(cohort["name"]),
            "title",
        )
        for step_idx, step in enumerate(cohort["main_flow"]):
            boxes[f"{prefix}_main_{step_idx}"] = Box(
                f"{prefix}_main_{step_idx}",
                layout["main_x"],
                y_main[step_idx],
                main_w,
                main_h,
                wrap_label(step["label"], step.get("n"), width=26),
                "main",
            )
        for exclusion in cohort.get("exclusions", []):
            step = int(exclusion["after_step"])
            if not 1 <= step < len(cohort["main_flow"]):
                continue
            y_mid = (y_main[step - 1] - main_h / 2 + y_main[step] + main_h / 2) / 2
            boxes[f"{prefix}_excl_{step}"] = Box(
                f"{prefix}_excl_{step}",
                layout["excl_x"],
                y_mid,
                excl_w,
                excl_h,
                wrap_label(exclusion["label"], exclusion.get("n"), width=29),
                "exclusion",
            )
        groups = cohort.get("final_groups", [])
        if groups:
            centers = outcome_centers(layout["main_x"], outcome_w, outcome_gap, len(groups))
            for group_idx, group in enumerate(groups):
                boxes[f"{prefix}_group_{group_idx}"] = Box(
                    f"{prefix}_group_{group_idx}",
                    centers[group_idx],
                    y_outcome,
                    outcome_w,
                    outcome_h,
                    wrap_label(group["label"], group.get("n"), width=20),
                    "outcome",
                )

    lines.extend(
        [
            "- PASS: fixed y positions used for main-flow boxes.",
            f"- PASS: outcome row y = {y_outcome:.3f} computed from final main-flow row and visible gap.",
        ]
    )
    overlaps: list[str] = []
    drawable = [box for box in boxes.values() if box.kind != "title"]
    for idx, first in enumerate(drawable):
        for second in drawable[idx + 1 :]:
            separated = (
                first.right <= second.left
                or second.right <= first.left
                or first.top <= second.bottom
                or second.top <= first.bottom
            )
            if not separated:
                overlaps.append(f"{first.key} overlaps {second.key}")
    if overlaps:
        lines.extend(f"- FAIL: {item}." for item in overlaps)
    else:
        lines.append("- PASS: no overlapping boxes detected.")
    lines.extend(validate_layout_geometry(config, boxes))
    return boxes, lines


def validate_layout_geometry(config: dict[str, Any], boxes: dict[str, Box]) -> list[str]:
    lines = ["## Connector Geometry"]
    for cohort_idx, cohort in enumerate(config["cohorts"]):
        prefix = f"cohort_{cohort_idx}"
        name = cohort.get("name", prefix)
        for step_idx in range(len(cohort["main_flow"]) - 1):
            upper = boxes[f"{prefix}_main_{step_idx}"]
            lower = boxes[f"{prefix}_main_{step_idx + 1}"]
            same_x = abs(upper.bottom_center[0] - lower.top_center[0]) <= LAYOUT_TOL
            ordered = upper.bottom > lower.top
            if same_x and ordered:
                gap = upper.bottom - lower.top
                lines.append(
                    f"- PASS: {name} main connector {step_idx + 1}->{step_idx + 2} "
                    f"uses bottom_center to top_center; visible gap = {gap:.3f}."
                )
            else:
                lines.append(
                    f"- FAIL: {name} main connector {step_idx + 1}->{step_idx + 2} "
                    "does not use aligned bottom/top anchors."
                )

        for exclusion in cohort.get("exclusions", []):
            step = int(exclusion["after_step"])
            key = f"{prefix}_excl_{step}"
            if key not in boxes:
                continue
            upper = boxes[f"{prefix}_main_{step - 1}"]
            lower = boxes[f"{prefix}_main_{step}"]
            exclusion_box = boxes[key]
            expected_y = (upper.bottom + lower.top) / 2
            y_ok = abs(exclusion_box.left_center[1] - expected_y) <= LAYOUT_TOL
            horizontal_ok = exclusion_box.left > upper.cx
            if y_ok and horizontal_ok:
                lines.append(
                    f"- PASS: {name} exclusion after step {step} branches from transition midpoint "
                    f"y = {expected_y:.3f} to exclusion left_center."
                )
            else:
                lines.append(
                    f"- FAIL: {name} exclusion after step {step} is not anchored to the transition midpoint."
                )

        if cohort.get("final_groups"):
            risk = boxes[f"{prefix}_main_{len(cohort['main_flow']) - 1}"]
            group_boxes = [boxes[f"{prefix}_group_{idx}"] for idx in range(len(cohort["final_groups"]))]
            split = outcome_split_geometry(risk, group_boxes)
            centers = [box.cx for box in group_boxes]
            gaps = [group_boxes[idx + 1].left - group_boxes[idx].right for idx in range(len(group_boxes) - 1)]
            span_center_ok = abs(((min(centers) + max(centers)) / 2) - risk.cx) <= LAYOUT_TOL
            equal_spacing_ok = all(abs(gap - gaps[0]) <= LAYOUT_TOL for gap in gaps) if gaps else True
            same_row_ok = all(abs(box.cy - group_boxes[0].cy) <= LAYOUT_TOL for box in group_boxes)
            vertical_ok = risk.bottom > split.junction[1] > max(box.top for box in group_boxes)
            endpoint_ok = all(
                abs(split.branch_junctions[idx][0] - group_boxes[idx].cx) <= LAYOUT_TOL
                and abs(split.child_ends[idx][0] - group_boxes[idx].top_center[0]) <= LAYOUT_TOL
                and abs(split.child_ends[idx][1] - group_boxes[idx].top_center[1]) <= LAYOUT_TOL
                for idx in range(len(group_boxes))
            )
            if span_center_ok and equal_spacing_ok and same_row_ok and vertical_ok and endpoint_ok:
                offsets = [box.cx - risk.cx for box in group_boxes]
                lines.append(
                    f"- PASS: {name} outcome split has {len(group_boxes)} final groups centered on main x = {risk.cx:.3f}; "
                    f"offsets = {[round(offset, 3) for offset in offsets]}; "
                    f"between-box gaps = {[round(gap, 3) for gap in gaps]}; "
                    f"central junction = ({split.junction[0]:.3f}, {split.junction[1]:.3f})."
                )
            else:
                lines.append(
                    f"- FAIL: {name} outcome split is not a centered trunk + explicit junction + child-branch structure."
                )
    return lines


def make_patch(box: Box) -> Rectangle:
    return Rectangle(
        (box.left, box.bottom),
        box.w,
        box.h,
        facecolor="white",
        edgecolor="#000000",
        linewidth=0.85,
        zorder=3,
    )


def draw_arrow(
    ax: Any,
    start: tuple[float, float],
    end: tuple[float, float],
    patchA: Any | None = None,
    patchB: Any | None = None,
    arrowstyle: str = "-|>",
    lw: float = 0.85,
) -> None:
    ax.add_patch(
        FancyArrowPatch(
            start,
            end,
            patchA=patchA,
            patchB=patchB,
            arrowstyle=arrowstyle,
            mutation_scale=9,
            linewidth=lw,
            color="#000000",
            shrinkA=0,
            shrinkB=0,
            connectionstyle="arc3,rad=0",
            zorder=1,
        )
    )


def draw_outcome_split(ax: Any, risk_box: Box, group_boxes: list[Box], patches: dict[str, Any]) -> None:
    split = outcome_split_geometry(risk_box, group_boxes)
    draw_arrow(
        ax,
        split.trunk_start,
        split.junction,
        patchA=patches[risk_box.key],
        arrowstyle="-",
    )
    for branch_junction in split.branch_junctions:
        if abs(branch_junction[0] - split.junction[0]) > LAYOUT_TOL:
            draw_arrow(ax, split.junction, branch_junction, arrowstyle="-")
    for group_box, branch_junction, child_end in zip(group_boxes, split.branch_junctions, split.child_ends):
        draw_arrow(ax, branch_junction, child_end, patchB=patches[group_box.key])


def draw_figure(config: dict[str, Any], boxes: dict[str, Box]) -> plt.Figure:
    fig_w, _ = compute_x_layouts(config)
    fig_h = 10
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    ax.set_xlim(0, fig_w)
    ax.set_ylim(1.4, 9.8)
    ax.axis("off")

    title = str(config.get("figure_title", "")).strip()
    if title:
        ax.text(fig_w / 2, 9.65, title, ha="center", va="center", fontsize=10.5, family="Arial")

    patches = {key: make_patch(box) for key, box in boxes.items() if box.kind != "title"}
    for patch in patches.values():
        ax.add_patch(patch)

    for cohort_idx, cohort in enumerate(config["cohorts"]):
        prefix = f"cohort_{cohort_idx}"
        title_box = boxes[f"{prefix}_title"]
        ax.text(title_box.cx, title_box.cy, title_box.text, ha="center", va="center", fontsize=10, family="Arial")
        for step_idx in range(len(cohort["main_flow"]) - 1):
            upper = boxes[f"{prefix}_main_{step_idx}"]
            lower = boxes[f"{prefix}_main_{step_idx + 1}"]
            draw_arrow(ax, upper.bottom_center, lower.top_center, patches[upper.key], patches[lower.key])
        for exclusion in cohort.get("exclusions", []):
            step = int(exclusion["after_step"])
            key = f"{prefix}_excl_{step}"
            if key not in boxes:
                continue
            upper = boxes[f"{prefix}_main_{step - 1}"]
            lower = boxes[f"{prefix}_main_{step}"]
            y_mid = (upper.bottom + lower.top) / 2
            draw_arrow(ax, (upper.cx, y_mid), boxes[key].left_center, patchB=patches[key])
        if cohort.get("final_groups"):
            risk = boxes[f"{prefix}_main_{len(cohort['main_flow']) - 1}"]
            group_boxes = [boxes[f"{prefix}_group_{idx}"] for idx in range(len(cohort["final_groups"]))]
            draw_outcome_split(ax, risk, group_boxes, patches)

    text_artists: list[tuple[Box, Any]] = []
    for box in boxes.values():
        if box.kind == "title":
            continue
        fontsize = 8.3 if box.kind != "exclusion" else 7.8
        artist = ax.text(
            box.cx,
            box.cy,
            box.text,
            ha="center",
            va="center",
            fontsize=fontsize,
            family="Arial",
            linespacing=1.12,
            color="#000000",
            zorder=4,
        )
        text_artists.append((box, artist))
    setattr(fig, "_flowchart_text_artists", text_artists)
    return fig


def validate_text_fit(fig: plt.Figure, padding_px: float = 2.0) -> list[str]:
    lines = ["## Text Fit QC"]
    fig.canvas.draw()
    renderer = fig.canvas.get_renderer()
    text_artists = getattr(fig, "_flowchart_text_artists", [])
    if not text_artists:
        lines.append("- FAIL: no flowchart text artists were available for text-fit QC.")
        return lines

    for box, artist in text_artists:
        text_bbox = artist.get_window_extent(renderer=renderer)
        ax = artist.axes
        lower_left = ax.transData.transform((box.left, box.bottom))
        upper_right = ax.transData.transform((box.right, box.top))
        box_left = min(lower_left[0], upper_right[0]) + padding_px
        box_right = max(lower_left[0], upper_right[0]) - padding_px
        box_bottom = min(lower_left[1], upper_right[1]) + padding_px
        box_top = max(lower_left[1], upper_right[1]) - padding_px
        if (
            text_bbox.x0 >= box_left
            and text_bbox.x1 <= box_right
            and text_bbox.y0 >= box_bottom
            and text_bbox.y1 <= box_top
        ):
            lines.append(f"- PASS: text fits inside {box.key}.")
        else:
            overflow = {
                "left_px": float(round(max(box_left - text_bbox.x0, 0), 1)),
                "right_px": float(round(max(text_bbox.x1 - box_right, 0), 1)),
                "bottom_px": float(round(max(box_bottom - text_bbox.y0, 0), 1)),
                "top_px": float(round(max(text_bbox.y1 - box_top, 0), 1)),
            }
            lines.append(f"- FAIL: text overflows {box.key}: {overflow}.")
    return lines


def add_ppt_arrowhead(connector: Any) -> None:
    line_properties = connector._element.spPr.ln
    if line_properties is None:
        return
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    line_properties.append(tail)


def ppt_xy(point: tuple[float, float], slide_h: float) -> tuple[Any, Any]:
    return Inches(point[0]), Inches(slide_h - point[1])


def ppt_connector(slide: Any, start: tuple[float, float], end: tuple[float, float], slide_h: float, arrow: bool = True) -> None:
    x1, y1 = ppt_xy(start, slide_h)
    x2, y2 = ppt_xy(end, slide_h)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = RGBColor(0, 0, 0)
    line.line.width = Pt(0.85)
    if arrow:
        add_ppt_arrowhead(line)


def add_ppt_text(shape: Any, text: str, font_size: float) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for line_idx, line_text in enumerate(text.splitlines() or [""]):
        p = frame.paragraphs[0] if line_idx == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line_text
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0, 0, 0)


def export_pptx(config: dict[str, Any], boxes: dict[str, Box], path: Path) -> str:
    if Presentation is None:
        return "PPTX export failed: python-pptx is unavailable."
    slide_w, _ = compute_x_layouts(config)
    slide_h = 10
    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = str(config.get("figure_title", "")).strip()
    if title:
        shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(slide_w - 1), Inches(0.35))
        add_ppt_text(shape, title, 10.5)

    for cohort_idx, cohort in enumerate(config["cohorts"]):
        prefix = f"cohort_{cohort_idx}"
        for step_idx in range(len(cohort["main_flow"]) - 1):
            upper = boxes[f"{prefix}_main_{step_idx}"]
            lower = boxes[f"{prefix}_main_{step_idx + 1}"]
            ppt_connector(slide, upper.bottom_center, lower.top_center, slide_h)
        for exclusion in cohort.get("exclusions", []):
            step = int(exclusion["after_step"])
            key = f"{prefix}_excl_{step}"
            if key not in boxes:
                continue
            upper = boxes[f"{prefix}_main_{step - 1}"]
            lower = boxes[f"{prefix}_main_{step}"]
            y_mid = (upper.bottom + lower.top) / 2
            ppt_connector(slide, (upper.cx, y_mid), boxes[key].left_center, slide_h)
        if cohort.get("final_groups"):
            risk = boxes[f"{prefix}_main_{len(cohort['main_flow']) - 1}"]
            group_boxes = [boxes[f"{prefix}_group_{idx}"] for idx in range(len(cohort["final_groups"]))]
            split = outcome_split_geometry(risk, group_boxes)
            ppt_connector(slide, split.trunk_start, split.junction, slide_h, arrow=False)
            for branch_junction in split.branch_junctions:
                if abs(branch_junction[0] - split.junction[0]) > LAYOUT_TOL:
                    ppt_connector(slide, split.junction, branch_junction, slide_h, arrow=False)
            for branch_junction, child_end in zip(split.branch_junctions, split.child_ends):
                ppt_connector(slide, branch_junction, child_end, slide_h)

    for box in boxes.values():
        if box.kind == "title":
            shape = slide.shapes.add_textbox(Inches(box.left), Inches(slide_h - box.top), Inches(box.w), Inches(box.h))
            add_ppt_text(shape, box.text, 10)
            continue
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(box.left),
            Inches(slide_h - box.top),
            Inches(box.w),
            Inches(box.h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(0, 0, 0)
        shape.line.width = Pt(0.85)
        add_ppt_text(shape, box.text, 8.2 if box.kind != "exclusion" else 7.7)

    prs.save(path)
    return "PPTX export passed."


def archive_existing(paths: list[Path], archive_dir: Path) -> None:
    archive_dir.mkdir(parents=True, exist_ok=True)
    for path in paths:
        if not path.exists():
            continue
        target = archive_dir / path.name
        counter = 1
        while target.exists():
            target = archive_dir / f"{path.stem}_archived_{counter}{path.suffix}"
            counter += 1
        shutil.move(str(path), str(target))


def write_qc_report(
    path: Path,
    count_ok: bool,
    layout_ok: bool,
    count_lines: list[str],
    layout_lines: list[str],
    export_paths: dict[str, Path],
    extra_lines: list[str],
    report_title: str = "Figure 1 Cohort Selection Flowchart QC Report",
) -> None:
    lines = [f"# {report_title}", ""]
    lines.extend(count_lines)
    lines.append("")
    lines.extend(layout_lines)
    lines.append("")
    lines.append("## Export Consistency")
    for label, file_path in export_paths.items():
        status = "PASS" if label == "qc" or file_path.exists() else "FAIL"
        lines.append(f"- {status}: {label}: `{file_path}`")
    lines.extend(f"- {line}" for line in extra_lines)
    lines.append("")
    lines.append("## Overall Status")
    if count_ok and layout_ok:
        lines.append("- PASS")
    elif not count_ok:
        lines.append("- FAIL: sample-size inconsistencies require correction.")
    else:
        lines.append("- FAIL: layout, connector geometry, or text-fit consistency requires correction.")
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a fixed-coordinate SCI cohort flowchart.")
    parser.add_argument("--config", type=Path, default=None, help="JSON or YAML config file.")
    parser.add_argument("--output-dir", type=Path, default=Path("outputs/figures"), help="Output directory.")
    parser.add_argument("--allow-qc-fail", action="store_true", help="Export draft even if count QC fails.")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    config = load_config(args.config)
    output_dir = args.output_dir
    output_dir.mkdir(parents=True, exist_ok=True)
    base = str(config.get("output_basename", "Figure_1_cohort_selection_flowchart"))
    report_title = str(config.get("qc_report_title", "Figure 1 Cohort Selection Flowchart QC Report"))
    paths = {
        "pdf": output_dir / f"{base}.pdf",
        "svg": output_dir / f"{base}.svg",
        "png": output_dir / f"{base}_600dpi.png",
        "pptx": output_dir / f"{base}_editable.pptx",
        "qc": output_dir / f"{base}_QC_report.md",
        "script": output_dir / "generate_scientific_cohort_flowchart.py",
    }
    archive_existing(list(paths.values()), output_dir / "archive")

    count_ok, count_lines = validate_counts(config)
    boxes, layout_lines = make_layout(config)
    pre_render_layout_ok = is_layout_pass(layout_lines)
    extra_lines: list[str] = []

    shutil.copyfile(Path(__file__), paths["script"])
    if (not count_ok or not pre_render_layout_ok) and not args.allow_qc_fail:
        write_qc_report(
            paths["qc"], count_ok, pre_render_layout_ok, count_lines, layout_lines, paths, extra_lines, report_title
        )
        print(f"QC failed. Report written to {paths['qc']}")
        return 2

    fig = draw_figure(config, boxes)
    layout_lines.extend(validate_text_fit(fig))
    layout_ok = is_layout_pass(layout_lines)

    if not layout_ok and not args.allow_qc_fail:
        write_qc_report(paths["qc"], count_ok, layout_ok, count_lines, layout_lines, paths, extra_lines, report_title)
        plt.close(fig)
        print(f"QC failed. Report written to {paths['qc']}")
        return 2

    fig.savefig(paths["pdf"], bbox_inches="tight", facecolor="white")
    fig.savefig(paths["svg"], bbox_inches="tight", facecolor="white")
    fig.savefig(paths["png"], dpi=600, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    extra_lines.append(export_pptx(config, boxes, paths["pptx"]))
    write_qc_report(paths["qc"], count_ok, layout_ok, count_lines, layout_lines, paths, extra_lines, report_title)
    print(f"Generated outputs in {output_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
